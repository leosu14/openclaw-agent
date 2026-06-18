from pptx import Presentation
from pptx.util import Inches
from themes.education_theme import *
from tools.image_tools import download_image


TYPE_LAYOUT = {

    "cover": 0,

    "warmup": 1,

    "image": 1,

    "vocab": 1,

    "reading": 1,

    "grammar": 1,

    "speaking": 1,

    "summary": 1,

    "homework": 1
}


def create_ppt(
    title,
    content,
    output
):

    prs = Presentation()

    blocks = content.split(
        "SLIDE:"
    )

    for block in blocks:

        if not block.strip():
            continue

        slide_type = "normal"
        slide_title = ""
        image_keyword = None

        bullets = []

        lines = (
            block
            .strip()
            .split("\n")
        )

        # ======================
        # Parse Block
        # ======================

        for line in lines:

            line = line.strip()

            if line.startswith("TYPE:"):

                slide_type = (
                    line
                    .replace(
                        "TYPE:",
                        ""
                    )
                    .strip()
                )

            elif line.startswith(
                "TITLE:"
            ):

                slide_title = (
                    line
                    .replace(
                        "TITLE:",
                        ""
                    )
                    .strip()
                )

            elif line.startswith(
                "IMAGE:"
            ):

                image_keyword = (
                    line
                    .replace(
                        "IMAGE:",
                        ""
                    )
                    .strip()
                )

            elif line.startswith("-"):

                bullets.append(
                    line
                )

        # ======================
        # Create Slide
        # ======================

        layout_id = TYPE_LAYOUT.get(
            slide_type,
            1
        )

        slide = prs.slides.add_slide(
            prs.slide_layouts[
                layout_id
            ]
        )

        # ======================
        # Title
        # ======================

        try:

            title_shape = slide.shapes.title

            title_shape.text = slide_title

            title_shape.text_frame.paragraphs[0].font.name = (
                FONT_NAME
            )

            title_shape.text_frame.paragraphs[0].font.size = (
                Pt(TITLE_SIZE)
            )

            title_shape.text_frame.paragraphs[0].font.color.rgb = (
                TITLE_COLOR
            )

        except:
            pass

        # ======================
        # Image
        # ======================

        if image_keyword:

            image_path = download_image(
                image_keyword
            )

            if image_path:

                try:

                    # Cover Page
                    if slide_type == "cover":

                        slide.shapes.add_picture(

                            image_path,

                            Inches(1),

                            Inches(1.3),

                            width=Inches(7)
                        )
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(
                            248,
                            250,
                            252
                        )

                    # Big Discussion Image
                    elif slide_type == "image":

                        slide.shapes.add_picture(

                            image_path,

                            Inches(0.5),

                            Inches(1.5),

                            width=Inches(6)
                        )

                    # Normal Image
                    else:

                        slide.shapes.add_picture(

                            image_path,

                            Inches(5),

                            Inches(1),

                            width=Inches(3)
                        )

                except Exception as e:

                    print(
                        f"Image insert error: {e}"
                    )

        # ======================
        # Content
        # ======================

        try:

            if len(
                slide.placeholders
            ) > 1:

                slide.placeholders[
                    1
                ].text = (
                    "\n".join(
                        bullets
                    )
                )

        except:
            pass

        # ======================
        # Summary Style
        # ======================

        if slide_type == "summary":

            try:

                slide.shapes.title.text = (
                    "🧠 Ideas Clave"
                )

                slide.shapes.title.text_frame.paragraphs[
                    0
                ].font.color.rgb = (
                    SUMMARY_COLOR
                )

            except:
                pass

        # ======================
        # Homework Style
        # ======================

        if slide_type == "homework":

            try:

                slide.shapes.title.text = (
                    "📝 Tarea"
                )

                slide.shapes.title.text_frame.paragraphs[
                    0
                ].font.color.rgb = (
                    HOMEWORK_COLOR
                )

            except:
                pass

        # ======================
        # Vocab Style
        # ======================

        if slide_type == "vocab":

            try:

                title_shape.text_frame.paragraphs[
                    0
                ].font.color.rgb = (
                    VOCAB_COLOR
                )

            except:
                pass



    prs.save(output)

    print(
        f"✅ PPT saved: {output}"
    )